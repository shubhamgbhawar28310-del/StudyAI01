import os
import re
import chardet
import traceback
import fitz  # PyMuPDF
import pdfplumber
import PyPDF2
import docx
import pptx
from io import BytesIO

class DocumentService:
    def __init__(self):
        # Configure maximum pages to process for different file types
        self.max_pdf_pages = 100
        self.max_pptx_slides = 100
        self.max_chunk_size = 500000  # Maximum characters per chunk (500KB of text)
    
    def extract_text_from_file(self, file_path, filename):
        """
        Extract text from file based on file type
        """
        try:
            file_extension = os.path.splitext(filename)[1].lower()
            
            if file_extension == '.pdf':
                return self.extract_pdf_text(file_path)
            elif file_extension == '.docx':
                return self.extract_docx_text(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                return self.extract_pptx_text(file_path)
            elif file_extension == '.txt':
                # For testing purposes, read text files directly
                return self._read_text_file(file_path)
            else:
                return f"File type {file_extension} not supported for text extraction. File name: {filename}"
        except Exception as e:
            print(f"Error extracting text from {filename}: {str(e)}")
            traceback.print_exc()
            raise Exception(f"Failed to extract text from file: {str(e)}")
    
    def extract_pdf_text(self, file_path):
        """
        Extract text from PDF file using multiple methods:
        1. PyMuPDF (fitz) - most reliable for various PDF types
        2. pdfplumber - for better layout preservation
        3. PyPDF2 - fallback method
        """
        try:
            # Try PyMuPDF first (most reliable)
            text = self._extract_with_pymupdf(file_path)
            if text and len(text.strip()) > 50:  # If we got meaningful text
                print(f"Successfully extracted text with PyMuPDF: {len(text)} characters")
                cleaned_text = self._clean_text(text)
                return self._improve_text_formatting(cleaned_text)
            
            # Try pdfplumber for better layout (slower but more accurate)
            text = self._extract_with_pdfplumber(file_path)
            if text and len(text.strip()) > 50:  # If we got meaningful text
                print(f"Successfully extracted text with pdfplumber: {len(text)} characters")
                cleaned_text = self._clean_text(text)
                return self._improve_text_formatting(cleaned_text)
            
            # Try PyPDF2 (faster but less accurate)
            text = self._extract_with_pypdf2(file_path)
            if text and len(text.strip()) > 50:  # If we got meaningful text
                print(f"Successfully extracted text with PyPDF2: {len(text)} characters")
                cleaned_text = self._clean_text(text)
                return self._improve_text_formatting(cleaned_text)
            
            # If still no good text, return what we have
            if text:
                cleaned_text = self._clean_text(text)
                formatted_text = self._improve_text_formatting(cleaned_text)
                return formatted_text + "\n\nNote: This document may contain scanned images. For better extraction, OCR processing would be needed."
            else:
                return "No text content found in PDF. This document may contain only images or be password protected."
                
        except Exception as e:
            print(f"PDF extraction error: {str(e)}")
            traceback.print_exc()
            raise Exception(f"Failed to parse PDF: {str(e)}")
    
    def _extract_with_pymupdf(self, file_path):
        """Extract text using PyMuPDF (fitz) - most reliable method"""
        try:
            text = ""
            doc = fitz.open(file_path)
            
            # Get total pages for progress tracking
            total_pages = len(doc)
            print(f"Processing PDF with PyMuPDF: {total_pages} pages")
            
            # Process all pages (limit to 100 for very large documents)
            pages_to_process = min(total_pages, self.max_pdf_pages)
            for page_num in range(pages_to_process):
                try:
                    page = doc.load_page(page_num)
                    page_text = page.get_text()
                    if page_text:
                        text += f"\n\n--- Page {page_num + 1} ---\n{page_text}"
                except Exception as page_error:
                    print(f"Error extracting page {page_num + 1} with PyMuPDF: {str(page_error)}")
                    text += f"\n\n--- Page {page_num + 1} ---\n[Error extracting text from this page]"
            
            doc.close()
            return text
        except Exception as e:
            print(f"PyMuPDF extraction failed: {str(e)}")
            return ""
    
    def _extract_with_pdfplumber(self, file_path):
        """Extract text using pdfplumber for better layout preservation"""
        try:
            text = ""
            with pdfplumber.open(file_path) as pdf:
                # Get total pages for progress tracking
                total_pages = len(pdf.pages)
                print(f"Processing PDF with pdfplumber: {total_pages} pages")
                
                # Process all pages (limit to 50 for very large documents)
                pages_to_process = min(total_pages, 50)
                for page_num in range(pages_to_process):
                    try:
                        page = pdf.pages[page_num]
                        page_text = page.extract_text()
                        if page_text:
                            text += f"\n\n--- Page {page_num + 1} ---\n{page_text}"
                    except Exception as page_error:
                        print(f"Error extracting page {page_num + 1} with pdfplumber: {str(page_error)}")
                        text += f"\n\n--- Page {page_num + 1} ---\n[Error extracting text from this page]"
            return text
        except Exception as e:
            print(f"pdfplumber extraction failed: {str(e)}")
            return ""
    
    def _extract_with_pypdf2(self, file_path):
        """Extract text using PyPDF2 with enhanced error handling"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                try:
                    pdf_reader = PyPDF2.PdfReader(file)
                    
                    # Check if PDF is encrypted
                    if pdf_reader.is_encrypted:
                        try:
                            # Try with empty password
                            pdf_reader.decrypt('')
                        except:
                            return "This PDF is password-protected and cannot be processed."
                    
                    # Get total pages for progress tracking
                    total_pages = len(pdf_reader.pages)
                    print(f"Processing PDF with PyPDF2: {total_pages} pages")
                    
                    # Process all pages (limit to 100 for very large documents)
                    pages_to_process = min(total_pages, self.max_pdf_pages)
                    for page_num in range(pages_to_process):
                        try:
                            page = pdf_reader.pages[page_num]
                            page_text = page.extract_text()
                            if page_text:
                                text += f"\n\n--- Page {page_num + 1} ---\n{page_text}"
                        except Exception as page_error:
                            print(f"Error extracting page {page_num + 1}: {str(page_error)}")
                            text += f"\n\n--- Page {page_num + 1} ---\n[Error extracting text from this page]"
                except Exception as pdf_error:
                    print(f"PyPDF2 reader error: {str(pdf_error)}")
                    return ""
            return text
        except Exception as e:
            print(f"PyPDF2 extraction failed: {str(e)}")
            return ""
    
    def extract_pptx_text(self, file_path):
        """
        Extract text from PPTX file preserving structure including:
        - Slide titles
        - Text boxes
        - Bullet points
        - Speaker notes
        - Tables and other elements
        """
        try:
            # Process PPTX file
            presentation = pptx.Presentation(file_path)
            text_content = []
            
            print(f"Processing presentation with {len(presentation.slides)} slides")
            
            for i, slide in enumerate(presentation.slides):
                try:
                    slide_text = []
                    
                    # Extract slide title (usually the first placeholder)
                    title = ""
                    if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title and hasattr(slide.shapes.title, 'text'):
                        title = slide.shapes.title.text.strip()
                    
                    # Extract content from all shapes
                    if hasattr(slide, 'shapes'):
                        for shape in slide.shapes:
                            try:
                                # Extract text from text frames
                                if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                                    # Skip the title shape as we already extracted it
                                    if hasattr(slide.shapes, 'title') and shape == slide.shapes.title:
                                        continue
                                        
                                    text_frame = shape.text_frame
                                    paragraphs = []
                                    for paragraph in text_frame.paragraphs:
                                        if paragraph.text.strip():
                                            # Simple approach: just add the text as-is
                                            # The bullet points should already be in the text
                                            paragraphs.append(paragraph.text.strip())
                                    
                                    if paragraphs:
                                        slide_text.extend(paragraphs)
                                
                                # Extract text from tables
                                if hasattr(shape, 'has_table') and shape.has_table:
                                    table = shape.table
                                    for row in table.rows:
                                        row_text = []
                                        for cell in row.cells:
                                            if cell.text_frame.text.strip():
                                                row_text.append(cell.text_frame.text.strip())
                                        if row_text:
                                            slide_text.append(" | ".join(row_text))
                            except Exception as shape_error:
                                print(f"Error processing shape in slide {i+1}: {str(shape_error)}")
                                continue
                    
                    # Extract speaker notes
                    notes_text = []
                    if hasattr(slide, 'has_notes_slide') and slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                        if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                            for paragraph in notes_slide.notes_text_frame.paragraphs:
                                if paragraph.text.strip():
                                    notes_text.append(f"NOTE: {paragraph.text.strip()}")
                    
                    # Format slide content
                    if title or slide_text or notes_text:
                        slide_content = []
                        if title:
                            slide_content.append(f"SLIDE {i+1}: {title}")
                        
                        if slide_text:
                            slide_content.append("CONTENT:")
                            for item in slide_text:
                                slide_content.append(f"  {item}")
                        
                        if notes_text:
                            slide_content.append("SPEAKER NOTES:")
                            for note in notes_text:
                                slide_content.append(f"  {note}")
                        
                        text_content.append("\n".join(slide_content))
                except Exception as slide_error:
                    print(f"Error processing slide {i+1}: {str(slide_error)}")
                    text_content.append(f"SLIDE {i+1}: [Error processing this slide]")
            
            if text_content:
                full_text = "\n\n".join(text_content)
                # Apply text formatting improvements
                cleaned_text = self._clean_text(full_text)
                return self._improve_text_formatting(cleaned_text)
            else:
                return f"No text content found in presentation."
                
        except Exception as e:
            print(f"PPTX extraction error: {str(e)}")
            traceback.print_exc()
            raise Exception(f"Failed to parse PPTX: {str(e)}")
    
    def extract_docx_text(self, file_path):
        """
        Extract text from DOCX file preserving structure
        """
        try:
            doc = docx.Document(file_path)
            text_content = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    # Add formatting indicators
                    if paragraph.style.name.startswith('Heading'):
                        text_content.append(f"\n{paragraph.style.name.upper()}: {paragraph.text.strip()}")
                    else:
                        text_content.append(paragraph.text.strip())
            
            # Extract tables
            for table_num, table in enumerate(doc.tables):
                text_content.append(f"\nTABLE {table_num + 1}:")
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_content.append(f"  {row_text}")
            
            full_text = "\n".join(text_content)
            cleaned_text = self._clean_text(full_text)
            return self._improve_text_formatting(cleaned_text)
                
        except Exception as e:
            print(f"DOCX extraction error: {str(e)}")
            traceback.print_exc()
            raise Exception(f"Failed to parse DOCX: {str(e)}")
    
    def _clean_text(self, text):
        """
        Clean extracted text by removing headers/footers, normalizing spaces, and fixing common OCR errors
        """
        if not text:
            return text
        
        # Remove excessive whitespace while preserving paragraph structure
        text = re.sub(r'[ \t]+', ' ', text)  # Normalize spaces and tabs
        text = re.sub(r'\n{3,}', '\n\n', text)  # Limit consecutive newlines to 2
        
        # Split into lines for processing
        lines = text.split('\n')
        cleaned_lines = []
        
        for line in lines:
            line = line.strip()
            # Skip empty lines
            if not line:
                continue
            # Skip lines that look like page numbers (standalone numbers)
            if re.match(r'^\d+$', line):
                continue
            # Skip lines that look like typical headers/footers
            if re.match(r'^\s*[A-Z][a-z]*\s+\d+\s*$', line):  # e.g., "Document 5"
                continue
            # Skip very short lines that might be artifacts (less than 3 characters)
            if len(line) < 3:
                continue
            cleaned_lines.append(line)
        
        # Join lines and normalize spacing
        cleaned_text = '\n'.join(cleaned_lines)
        
        # Final cleanup
        cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)  # Limit consecutive newlines
        
        return cleaned_text.strip()
    
    def _improve_text_formatting(self, text):
        """
        Improve text formatting for better readability and AI processing
        """
        if not text:
            return text
            
        # Clean up multiple spaces
        text = re.sub(r' {2,}', ' ', text)
        
        # Ensure proper spacing around punctuation
        text = re.sub(r'\s+([,.!?;:])', r'\1', text)
        text = re.sub(r'([,.!?;:])([A-Za-z])', r'\1 \2', text)
        
        # Fix bullet point formatting - remove extra newlines before bullets
        text = re.sub(r'\n\s*\n(\s*•)', r'\n\1', text)
        
        # Fix numbered list formatting
        text = re.sub(r'\n\s*\n(\s*\d+\.)', r'\n\1', text)
        
        # Remove extra whitespace at the beginning of lines
        lines = text.split('\n')
        cleaned_lines = [line.lstrip() for line in lines]
        text = '\n'.join(cleaned_lines)
        
        return text.strip()

    def chunk_text(self, text):
        """
        Split large text into manageable chunks while preserving paragraph structure
        Returns a list of text chunks
        """
        if not text or len(text) <= self.max_chunk_size:
            return [text] if text else []
            
        chunks = []
        paragraphs = re.split(r'\n\s*\n', text)
        current_chunk = ""
        
        for paragraph in paragraphs:
            # If adding this paragraph would exceed chunk size, start a new chunk
            if len(current_chunk) + len(paragraph) > self.max_chunk_size:
                # If current paragraph alone exceeds chunk size, split it further
                if len(paragraph) > self.max_chunk_size:
                    words = paragraph.split()
                    temp_para = ""
                    
                    for word in words:
                        if len(temp_para) + len(word) + 1 > self.max_chunk_size:
                            if current_chunk:
                                chunks.append(current_chunk.strip())
                            chunks.append(temp_para.strip())
                            temp_para = word
                            current_chunk = ""
                        else:
                            temp_para += " " + word if temp_para else word
                    
                    current_chunk = temp_para
                else:
                    # Store current chunk and start new one with this paragraph
                    chunks.append(current_chunk.strip())
                    current_chunk = paragraph
            else:
                # Add paragraph separator if needed
                if current_chunk:
                    current_chunk += "\n\n"
                current_chunk += paragraph
        
        # Add the last chunk if it's not empty
        if current_chunk:
            chunks.append(current_chunk.strip())
            
        return chunks
    
    def analyze_text(self, content, title="Document"):
        """
        Analyze text content (this would typically call an AI service)
        """
        return {
            "title": title,
            "wordCount": len(content.split()),
            "characterCount": len(content),
            "summary": content[:200] + ("..." if len(content) > 200 else ""),
            "mainTopics": ["Topic 1", "Topic 2", "Topic 3"]  # This would be AI-generated
        }
    
    def _read_text_file(self, file_path):
        """Read text file with automatic encoding detection"""
        try:
            # First try to detect encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                encoding = chardet.detect(raw_data)['encoding']
            
            # Then read with detected encoding
            with open(file_path, 'r', encoding=encoding) as f:
                return f.read()
        except Exception as e:
            # Fallback to common encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except:
                    continue
            # If all fail, raise the original exception
            raise e